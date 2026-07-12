#!/usr/bin/env python3
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parents[1]
ART = ROOT / "docs/demos/artifacts"
OUT = ART / "value-mechanism-evidence-deck.pptx"

W, H = 13.333, 7.5
BG = "081521"
PANEL = "102436"
PANEL2 = "152C40"
INK = "F4F8FC"
MUTED = "AFC2D2"
BLUE = "7FB7FF"
MINT = "64E6C5"
PEACH = "FFB87A"
PINK = "EE8FB7"
LINE = "29475E"
WHITE = "FFFFFF"
DARK = "07131F"


def rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_fill(shape, color, transparency=0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.fill.transparency = transparency


def set_line(shape, color, width=1):
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width)


def box(slide, x, y, w, h, fill=PANEL, line=LINE, radius=True):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(s, fill)
    set_line(s, line, 1)
    return s


def textbox(slide, x, y, w, h, text, size: float = 18, color=INK, bold=False,
            font="DejaVu Sans", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
            margin=0.04, fit=False):
    s = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = s.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)
    if fit:
        tf.fit_text(font_family=font, max_size=size)
    return s


def rich_text(slide, x, y, w, h, runs, size: float = 18, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    s = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = s.text_frame
    tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(.04)
    tf.margin_top = tf.margin_bottom = Inches(.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.alignment = align
    for text, color, bold in runs:
        r = p.add_run(); r.text = text
        r.font.name = "DejaVu Sans"; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = rgb(color)
    return s


def label(slide, text, x=.62, y=.38, color=MINT):
    textbox(slide, x, y, 5.6, .28, text.upper(), 9.5, color, True)


def footer(slide, n):
    textbox(slide, .62, 7.12, 5.6, .2, "VALUE–ACTION MECHANISTIC ANALYSIS", 8, MUTED, True)
    textbox(slide, 12.1, 7.1, .55, .22, f"0{n}", 8, MUTED, True, align=PP_ALIGN.RIGHT)


def add_bg(slide):
    bg = slide.background.fill
    bg.solid(); bg.fore_color.rgb = rgb(BG)
    # subtle decorative circles
    for x, y, w, c, trans in [(10.6,-1.0,3.8,BLUE,82),(-1.1,5.6,3.2,MINT,88)]:
        s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(w))
        set_fill(s, c, trans); s.line.fill.background()


def add_title(slide, title, subtitle=None, y=.83):
    textbox(slide, .62, y, 12.05, .72, title, 28, INK, True)
    if subtitle:
        textbox(slide, .65, y+.78, 11.3, .52, subtitle, 12.5, MUTED)


prs = Presentation()
prs.slide_width = Inches(W); prs.slide_height = Inches(H)
blank = prs.slide_layouts[6]

# Slide 1 — title
s = prs.slides.add_slide(blank); add_bg(s); label(s, "Research methods · mechanistic interpretability")
textbox(s, .72, 1.35, 11.1, 1.25, "From behavioral association\nto a mechanistic claim", 34, INK, True)
textbox(s, .76, 2.83, 9.6, .65, "An evidence standard for deciding whether a sparse feature participates in value-guided choice.", 16, MUTED)
# visual line
for i, (name, c) in enumerate([("DIRECTION", BLUE), ("LOCATION", MINT), ("TIMING", PEACH), ("INTERVENTION", PINK)]):
    x = .76 + i*2.92
    sh = box(s, x, 4.28, 2.55, 1.22, PANEL, c)
    textbox(s, x+.18, 4.49, 2.18, .25, f"0{i+1}", 10, c, True)
    textbox(s, x+.18, 4.85, 2.18, .3, name, 13, INK, True)
    if i < 3:
        textbox(s, x+2.58, 4.72, .3, .25, "→", 17, MUTED, True, align=PP_ALIGN.CENTER)
rich_text(s, 2.25, 6.18, 9.8, .35, [("CORE PRINCIPLE  ", MINT, True), ("A feature earns a stronger name only when independent evidence agrees.", INK, False)], 13)
footer(s, 1)

# Slide 2 — hero evidence ladder
s = prs.slides.add_slide(blank); add_bg(s); label(s, "Evidence ladder")
add_title(s, "What would count as a value mechanism?", "A directional correlation is the beginning of the analysis—not the conclusion.")
formula = box(s, .66, 1.9, 12.0, .65, "0D2030", LINE)
textbox(s, .9, 2.04, 4.0, .28, "SIGNED SUPPORT", 9, BLUE, True)
textbox(s, 4.1, 2.01, 5.0, .34, "D_f(x) = s(x)[Δ_f,A(x) − Δ_f,B(x)]", 15, INK, True, font="DejaVu Serif", align=PP_ALIGN.CENTER)
textbox(s, 10.0, 2.04, 2.35, .28, "D_f(x) > 0", 11, MUTED, True, align=PP_ALIGN.RIGHT)
steps = [
    ("01", "DIRECTION", "The signed contrast follows the value-side action across option order.", BLUE),
    ("02", "LOCATION", "Activity tracks the relevant option span rather than prompt scaffolding.", MINT),
    ("03", "TIMING", "Evidence appears before the answer, where it could influence the decision.", PEACH),
    ("04", "INTERVENTION", "Ablation or patching changes the choice under matched controls.", PINK),
]
for i,(num,title,desc,c) in enumerate(steps):
    x=.66+i*3.02
    box(s,x,2.88,2.72,2.48,PANEL,c)
    textbox(s,x+.2,3.08,.72,.38,num,16,c,True)
    textbox(s,x+.2,3.59,2.3,.35,title,12,INK,True)
    textbox(s,x+.2,4.08,2.28,.86,desc,11,MUTED)
    if i<3:
        circ=s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+2.69), Inches(3.84), Inches(.36), Inches(.36))
        set_fill(circ,DARK);set_line(circ,LINE,1)
        textbox(s,x+2.71,3.875,.31,.2,"→",10,MUTED,True,align=PP_ALIGN.CENTER)
claim=box(s,.66,5.72,12.0,.93,"123043",MINT)
rich_text(s,.92,5.95,11.45,.42,[("CLAIM RULE  ",MINT,True),("Call it a candidate until direction, location, timing, and intervention converge.",INK,False)],14,valign=MSO_ANCHOR.MIDDLE)
footer(s, 2)

# Slide 3 — confounds
s = prs.slides.add_slide(blank); add_bg(s); label(s, "Confound analysis", color=PEACH)
add_title(s, "One activation gap, three competing stories", "A ranked feature is ambiguous until competing explanations are tested.")
left=box(s,.66,1.92,4.0,4.72,"102436",BLUE)
textbox(s,.94,2.18,3.45,.28,"THE TEMPTING INTERPRETATION",9,BLUE,True)
textbox(s,.94,2.68,3.15,.88,"“This feature\nrepresents the value.”",23,INK,True)
textbox(s,.94,3.88,3.2,.78,"Observed: activation is larger when the model chooses the value-consistent option.",12,MUTED)
textbox(s,.94,5.22,3.2,.62,"Useful evidence—\nbut not yet identity or causality.",13,PEACH,True)
confounds=[
("PROMPT SCAFFOLD", "The feature tracks answer format, option order, or repeated instructions.", BLUE),
("CONTENT MATCH", "The feature recognizes words in one option without participating in choice.", MINT),
("DECISION CONSEQUENCE", "The feature activates after the answer direction is already determined.", PINK),
]
for i,(t,d,c) in enumerate(confounds):
    y=1.92+i*1.58
    box(s,5.0,y,7.66,1.34,PANEL2,c)
    textbox(s,5.28,y+.22,1.58,.25,f"0{i+1}",10,c,True)
    textbox(s,6.52,y+.2,2.0,.3,t,11,INK,True)
    textbox(s,8.58,y+.18,3.7,.72,d,11,MUTED)
footer(s, 3)

# Slide 4 — reporting standard
s = prs.slides.add_slide(blank); add_bg(s); label(s, "Claim discipline", color=PINK)
add_title(s, "Match the claim to the evidence", "The vocabulary should change as the evidence stack becomes more demanding.")
tiers=[
("DESCRIPTIVE", "associated feature", "Signed activation or contribution differs across choices.", BLUE, .66, 2.05, 3.62),
("LOCALIZED", "decision-relevant candidate", "Direction survives controls and appears at the relevant span or pre-answer site.", MINT, 4.85, 2.05, 3.62),
("CAUSAL", "mechanistic contributor", "Matched intervention changes the decision in the predicted direction.", PINK, 9.04, 2.05, 3.62),
]
for tier,name,desc,c,x,y,w in tiers:
    box(s,x,y,w,3.58,PANEL,c)
    textbox(s,x+.23,y+.25,w-.46,.25,tier,9,c,True)
    textbox(s,x+.23,y+.79,w-.46,.72,name,19,INK,True)
    textbox(s,x+.23,y+1.75,w-.46,.9,desc,11.5,MUTED)
    textbox(s,x+.23,y+2.93,w-.46,.25,"claim ceiling",9,c,True)
for x in [4.42,8.61]:
    textbox(s,x,3.64,.38,.3,"→",18,MUTED,True,align=PP_ALIGN.CENTER)
box(s,.66,6.02,12.0,.58,"0D2030",LINE)
textbox(s,.91,6.17,11.5,.23,"Mixed evidence is a result: preserve disagreement instead of averaging it into certainty.",12.5,INK,True,align=PP_ALIGN.CENTER)
footer(s, 4)

ART.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(OUT)
